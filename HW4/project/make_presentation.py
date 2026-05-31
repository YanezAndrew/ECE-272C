"""Generate Presentation.pptx for HW4."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x11, 0x17)   # dark background
ACCENT   = RGBColor(0x63, 0x66, 0xF1)   # indigo
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xE0, 0xE0, 0xE0)
MUTED    = RGBColor(0x6B, 0x72, 0x80)
GREEN    = RGBColor(0x86, 0xEF, 0xAC)
YELLOW   = RGBColor(0xFD, 0xE6, 0x8A)
RED      = RGBColor(0xFC, 0xA5, 0xA5)
DARKCARD = RGBColor(0x1E, 0x21, 0x30)

W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, left, top, width, height,
                   title=None, title_size=14, item_size=13,
                   bg_color=DARKCARD, title_color=ACCENT, item_color=LIGHT):
    add_rect(slide, left, top, width, height, bg_color)
    y = top + Inches(0.15)
    if title:
        add_text(slide, title, left + Inches(0.2), y,
                 width - Inches(0.4), Inches(0.35),
                 size=title_size, bold=True, color=title_color)
        y += Inches(0.35)
    for item in items:
        add_text(slide, f"• {item}", left + Inches(0.2), y,
                 width - Inches(0.4), Inches(0.32),
                 size=item_size, color=item_color)
        y += Inches(0.32)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s)
    # accent bar top
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    # title
    add_text(s, "Orchestrated AI Analytics System",
             Inches(1), Inches(1.8), Inches(11.33), Inches(1.1),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # subtitle
    add_text(s, "ECE 157C / ECE 272C — Homework 4",
             Inches(1), Inches(3.0), Inches(11.33), Inches(0.6),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, "200+ Financial Indicators of US Stocks (2014–2018)",
             Inches(1), Inches(3.7), Inches(11.33), Inches(0.5),
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, "Andrew Yanez · May 31, 2026",
             Inches(1), Inches(4.4), Inches(11.33), Inches(0.4),
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
    # accent bar bottom
    add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), ACCENT)


def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "System Architecture", Inches(0.5), Inches(0.2),
             Inches(8), Inches(0.6), size=28, bold=True, color=WHITE)

    components = [
        ("Orchestration Node", "Routes every question to analytics pipeline\nor DuckDuckGo web search. Drives validator\nretry loop (max 3 attempts)."),
        ("Enhanced Analytics Agent", "Autonomously inspects datasets, generates\n& executes Python in a persistent sandbox,\ndecides when the analysis is complete."),
        ("HW3 Validator Agent", "Independent semantic checker — reasons\nabout correctness & completeness, returns\nPASS / RETRY / SUSPICIOUS."),
        ("Web Search", "DuckDuckGo via LangChain for generic\nquestions. Returns grounded answer\nwith source citations."),
    ]

    xs = [Inches(0.3), Inches(3.55), Inches(6.8), Inches(10.05)]
    for i, (title, desc) in enumerate(components):
        x = xs[i]
        add_rect(s, x, Inches(1.1), Inches(3.0), Inches(5.6), DARKCARD)
        add_rect(s, x, Inches(1.1), Inches(3.0), Inches(0.06), ACCENT)
        add_text(s, title, x + Inches(0.15), Inches(1.25),
                 Inches(2.7), Inches(0.5), size=13, bold=True, color=ACCENT)
        add_text(s, desc, x + Inches(0.15), Inches(1.85),
                 Inches(2.7), Inches(2.5), size=11, color=LIGHT)

    # flow arrows label
    add_text(s, "User Question → Orchestrator → Analytics Agent ↔ Validator  |  or  |  Orchestrator → Web Search → Final Answer",
             Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4),
             size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_dataset(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Dataset Overview", Inches(0.5), Inches(0.2),
             Inches(8), Inches(0.6), size=28, bold=True, color=WHITE)

    stats = [
        ("5 Years", "2014 – 2018"),
        ("~4,400", "avg stocks/year"),
        ("225", "columns per file"),
        ("11", "sectors covered"),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.3 + i * 3.25)
        add_rect(s, x, Inches(1.0), Inches(3.0), Inches(1.5), DARKCARD)
        add_text(s, val,   x + Inches(0.15), Inches(1.1),
                 Inches(2.7), Inches(0.7), size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, label, x + Inches(0.15), Inches(1.75),
                 Inches(2.7), Inches(0.4), size=12, color=MUTED, align=PP_ALIGN.CENTER)

    col_groups = [
        ("Income Statement", ["Revenue", "Net Income", "EPS", "Gross Profit", "EBITDA"]),
        ("Balance Sheet",    ["Total Assets", "Total Debt", "Cash", "Shareholders Equity"]),
        ("Profitability",    ["Gross Margin", "Net Profit Margin", "ROE", "ROA", "ROIC"]),
        ("Valuation",        ["PE Ratio", "PB Ratio", "EV/EBITDA", "Debt to Equity"]),
        ("Growth",           ["Revenue Growth", "EPS Growth", "3Y/5Y/10Y Series"]),
        ("Target",           ["Class (1=outperformed next yr)", "Next-year price var %"]),
    ]
    xs2 = [Inches(0.3), Inches(4.55), Inches(8.8)]
    ys2 = [Inches(2.8), Inches(4.9)]
    for i, (title, items) in enumerate(col_groups):
        x = xs2[i % 3]
        y = ys2[i // 3]
        add_bullet_box(s, items, x, y, Inches(4.0), Inches(1.8),
                       title=title, title_size=12, item_size=10)


def slide_scenario_overview(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Analytics Scenario: Technology Sector Deep Dive (2014–2016)",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=24, bold=True, color=WHITE)
    add_text(s, "5-turn continuous conversation exploring Technology sector performance across three years",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
             size=13, color=MUTED)

    turns = [
        ("Turn 1", "Top 10 Tech stocks by\nmarket cap in 2014",
         "PASS", "Apple $599B, Google $360B,\nMicrosoft $344B led the sector", GREEN),
        ("Turn 2", "Average revenue trend of\nthose stocks (2014–2016)",
         "PASS", "$3.89B → $4.10B → $4.23B\n+$345M total growth", GREEN),
        ("Turn 3", "Net profit margin:\nTech vs Healthcare 2014–2016",
         "SUSPICIOUS", "Chart/text inconsistency flagged\nby validator (column mismatch)", YELLOW),
        ("Turn 4", "Highest EPS Growth in\nTechnology sector (2016)",
         "PASS", "IEC 100%, SPSC 68.5%,\nDIOD 53.25%", GREEN),
        ("Turn 5", "Most undervalued Tech stocks\n(PE, ROE, Rev Growth 2016)",
         "SUSPICIOUS", "PE=0 stocks flagged as\nmisleading by validator", YELLOW),
    ]

    for i, (turn, question, verdict, finding, vcolor) in enumerate(turns):
        x = Inches(0.3 + i * 2.6)
        add_rect(s, x, Inches(1.4), Inches(2.45), Inches(5.7), DARKCARD)
        add_rect(s, x, Inches(1.4), Inches(2.45), Inches(0.06), ACCENT)
        add_text(s, turn, x + Inches(0.12), Inches(1.5),
                 Inches(2.2), Inches(0.3), size=11, bold=True, color=ACCENT)
        add_text(s, question, x + Inches(0.12), Inches(1.85),
                 Inches(2.2), Inches(0.9), size=10, color=LIGHT)
        # verdict badge
        add_rect(s, x + Inches(0.12), Inches(2.85), Inches(2.2), Inches(0.3), BG)
        add_text(s, f"Validator: {verdict}", x + Inches(0.12), Inches(2.85),
                 Inches(2.2), Inches(0.3), size=10, bold=True, color=vcolor)
        add_text(s, finding, x + Inches(0.12), Inches(3.25),
                 Inches(2.2), Inches(1.5), size=9, color=MUTED)


def slide_key_finding_1(prs):
    """Tech sector market cap & revenue growth."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Finding 1: Technology Sector Revenue Growth (2014–2016)",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=26, bold=True, color=WHITE)

    # big numbers
    data = [("2014", "$3.89B"), ("2015", "$4.10B"), ("2016", "$4.23B")]
    for i, (year, val) in enumerate(data):
        x = Inches(0.5 + i * 2.8)
        add_rect(s, x, Inches(1.0), Inches(2.5), Inches(1.8), DARKCARD)
        add_text(s, year, x + Inches(0.1), Inches(1.1),
                 Inches(2.3), Inches(0.4), size=14, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(s, val,  x + Inches(0.1), Inches(1.5),
                 Inches(2.3), Inches(0.6), size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(s, "↑ +8.7% total growth over two years",
             Inches(0.5), Inches(3.0), Inches(8.0), Inches(0.4),
             size=14, bold=True, color=GREEN)

    add_bullet_box(s, [
        "Top 10 Tech stocks (by 2014 market cap) grew average revenue consistently year-over-year",
        "Apple, Google, and Microsoft accounted for the majority of aggregate revenue",
        "Revenue growth was steady even as broader market conditions shifted",
        "Cross-year analysis performed by joining 2014, 2015, and 2016 CSV files on ticker index",
    ], Inches(0.5), Inches(3.55), Inches(8.0), Inches(3.1),
       title="Key Observations", title_size=13, item_size=12)

    add_bullet_box(s, [
        "Top 10 by 2014 Market Cap:",
        "AAPL · GOOGL · MSFT",
        "FB · TSM · ORCL",
        "INTC · IBM · AMZN · SAP",
    ], Inches(8.8), Inches(1.0), Inches(4.2), Inches(5.7),
       title="Stocks Analyzed", title_size=12, item_size=11)


def slide_key_finding_2(prs):
    """EPS Growth leaders & validator behavior."""
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Finding 2: Technology Sector EPS Growth Leaders (2016)",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=26, bold=True, color=WHITE)

    stocks = [
        ("IEC",  "100.00%", "1st"),
        ("SPSC", " 68.50%", "2nd"),
        ("DIOD", " 53.25%", "3rd"),
    ]
    colors = [ACCENT, RGBColor(0x81, 0x8C, 0xF8), RGBColor(0xA5, 0xB4, 0xFC)]
    for i, (ticker, pct, rank) in enumerate(stocks):
        x = Inches(0.5 + i * 4.2)
        add_rect(s, x, Inches(1.0), Inches(3.8), Inches(2.2), DARKCARD)
        add_text(s, rank,   x + Inches(0.15), Inches(1.1),
                 Inches(3.5), Inches(0.35), size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(s, ticker, x + Inches(0.15), Inches(1.45),
                 Inches(3.5), Inches(0.55), size=30, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        add_text(s, pct,    x + Inches(0.15), Inches(2.0),
                 Inches(3.5), Inches(0.55), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, "EPS Growth", x + Inches(0.15), Inches(2.55),
                 Inches(3.5), Inches(0.3), size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_bullet_box(s, [
        "Agent loaded 2016_Financial_Data.csv, filtered Sector == 'Technology', sorted by EPS Growth descending",
        "Validator initially returned RETRY — requested explicit sector-filter confirmation in the result",
        "On retry, agent included sector count and column source in result — validator returned PASS",
        "Demonstrates validator improving result quality, not just checking execution success",
        "IEC's 100% EPS growth likely reflects recovery from a low/negative prior-year baseline",
    ], Inches(0.5), Inches(3.4), Inches(12.5), Inches(3.7),
       title="Analysis Details & Validator Behavior", title_size=13, item_size=12)


def slide_validator(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Validator Agent: Catching Real Errors",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=26, bold=True, color=WHITE)

    # Case 1
    add_rect(s, Inches(0.3), Inches(1.0), Inches(6.1), Inches(5.8), DARKCARD)
    add_rect(s, Inches(0.3), Inches(1.0), Inches(6.1), Inches(0.06), YELLOW)
    add_text(s, "SUSPICIOUS: Chart/Text Inconsistency",
             Inches(0.45), Inches(1.1), Inches(5.8), Inches(0.4),
             size=13, bold=True, color=YELLOW)
    add_text(s, "Question: Compare net profit margin of Tech vs Healthcare (2014–2016)",
             Inches(0.45), Inches(1.6), Inches(5.8), Inches(0.5),
             size=11, color=LIGHT)
    add_bullet_box(s, [
        "Text answer: Healthcare −35.05%, Tech −1.07% (negative)",
        "Generated charts: both sectors showed ~+1.0 (positive!)",
        "Root cause: agent used 'Net Profit Margin' for text but",
        "  'netProfitMargin' (different column, ~1.0) for chart",
        "Validator flagged internal contradiction correctly",
        "Retried 3× — column ambiguity persisted across all attempts",
    ], Inches(0.45), Inches(2.2), Inches(5.8), Inches(4.3),
       title="What happened", title_size=11, item_size=10, bg_color=BG)

    # Case 2
    add_rect(s, Inches(6.9), Inches(1.0), Inches(6.1), Inches(5.8), DARKCARD)
    add_rect(s, Inches(6.9), Inches(1.0), Inches(6.1), Inches(0.06), YELLOW)
    add_text(s, "RETRY: Incomplete Multi-Step Analysis",
             Inches(7.05), Inches(1.1), Inches(5.8), Inches(0.4),
             size=13, bold=True, color=YELLOW)
    add_text(s, "Question: Top 10 revenue-growth stocks in 2015 → how many Class=1 in 2016?",
             Inches(7.05), Inches(1.6), Inches(5.8), Inches(0.5),
             size=11, color=LIGHT)
    add_bullet_box(s, [
        "Agent answered '5 out of 10 outperformed' (50%)",
        "Validator: no ticker list shown, no join result table,",
        "  no visualization, 50% figure unsubstantiated",
        "Retried 3× — agent continued summarizing without",
        "  surfacing the underlying cross-year join data",
        "Final output accepted after max retries exhausted",
    ], Inches(7.05), Inches(2.2), Inches(5.8), Inches(4.3),
       title="What happened", title_size=11, item_size=10, bg_color=BG)


def slide_web_search(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Web Search: Generic-Domain Questions",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=26, bold=True, color=WHITE)

    examples = [
        (
            "What caused the 2016 oil-price crash?",
            "The crash was driven by a near-doubling of U.S. oil production from 2008 levels "
            "due to shale fracking technology improvements. This created a global oversupply, "
            "collapsed prices, and diminished OPEC's ability to control markets.",
            ["Wikipedia – 2010s oil glut", "Wikipedia – 2014–2016 oil market chronology",
             "Investopedia", "Wikipedia – Price of oil"],
        ),
        (
            "What is NVIDIA's primary business segment?",
            "NVIDIA's primary segment is Data Center, which has become the dominant revenue "
            "driver due to explosive demand for AI infrastructure from hyperscalers including "
            "Microsoft, Amazon, Google, and Meta.",
            ["MarketScreener", "StockDividendScreener",
             "Investopedia – How Nvidia Makes Money", "FourWeekMBA", "PitchGrade"],
        ),
    ]

    for i, (q, a, cites) in enumerate(examples):
        y = Inches(1.1 + i * 3.1)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(2.8), DARKCARD)
        add_text(s, f"Q: {q}", Inches(0.5), y + Inches(0.15),
                 Inches(12.3), Inches(0.4), size=13, bold=True, color=ACCENT)
        add_text(s, a, Inches(0.5), y + Inches(0.6),
                 Inches(8.5), Inches(1.4), size=11, color=LIGHT)
        cite_text = "  Sources: " + " · ".join(cites)
        add_text(s, cite_text, Inches(0.5), y + Inches(2.1),
                 Inches(12.3), Inches(0.4), size=9, color=MUTED)

    add_text(s, "Routing: Orchestrator classified both as 'web' → DuckDuckGo retrieved 5 results → GPT-4o synthesized grounded answer → citations returned to UI",
             Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
             size=9, color=MUTED, align=PP_ALIGN.CENTER)


def slide_lessons(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Lessons Learned & Future Work",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=26, bold=True, color=WHITE)

    add_bullet_box(s, [
        "Providing all 225 column names in the prompt was essential — early tests with only 20 columns caused systematic hallucinations",
        "Validator's SUSPICIOUS verdict for chart/text inconsistency proved more valuable than simple execution-success checks",
        "Stopping condition evaluation must be separate from code generation to avoid conflating 'code ran' with 'question answered'",
        "Cross-year joins produce _x/_y column suffixes — LLMs need explicit guidance to handle this reliably",
        "PE=0 stocks are unprofitable, not undervalued — domain-specific checks in the validator prompt are needed",
    ], Inches(0.3), Inches(1.0), Inches(12.7), Inches(3.0),
       title="Lessons Learned", title_size=14, item_size=12)

    add_bullet_box(s, [
        "Add a data-cleaning pre-step that resolves duplicate column names before analysis begins",
        "Provide cross-year join templates in the analytics agent prompt to handle _x/_y suffixes",
        "Add confidence scoring to validator verdicts — surface SUSPICIOUS as a warning rather than always retrying",
        "Stream intermediate execution results to the frontend so users see iteration progress in real time",
        "Support natural language SQL queries across yearly tables for more reliable multi-table operations",
    ], Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.9),
       title="Future Improvements", title_size=14, item_size=12)


def slide_summary(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, W, Inches(0.08), ACCENT)
    add_text(s, "Summary", Inches(0.5), Inches(0.2),
             Inches(12), Inches(0.6), size=28, bold=True, color=WHITE)

    items = [
        ("Orchestration", "GPT-4o classifier routes every question to analytics pipeline or DuckDuckGo web search"),
        ("Iterative Analytics", "Persistent Python sandbox with up to 5 reasoning/execution iterations per question"),
        ("Semantic Validation", "Independent validator reasons about correctness, completeness, and numerical consistency — not just execution success"),
        ("Temporal Analysis", "5-year dataset (2014–2018) enables cross-year joins, trend detection, and sector comparisons"),
        ("System Design", "FastAPI backend + SSE streaming + chatbot UI extended from HW2 with year filters, validator badges, citations, and agent traces"),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(1.1 + i * 1.18)
        add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.05), DARKCARD)
        add_text(s, title, Inches(0.5), y + Inches(0.1),
                 Inches(3.0), Inches(0.4), size=13, bold=True, color=ACCENT)
        add_text(s, desc, Inches(3.6), y + Inches(0.1),
                 Inches(9.2), Inches(0.75), size=12, color=LIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_title(prs)
    slide_architecture(prs)
    slide_dataset(prs)
    slide_scenario_overview(prs)
    slide_key_finding_1(prs)
    slide_key_finding_2(prs)
    slide_validator(prs)
    slide_web_search(prs)
    slide_lessons(prs)
    slide_summary(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
